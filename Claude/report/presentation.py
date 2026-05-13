"""
reports/presentation.py
=======================
Generazione report PowerPoint per i risultati di training.

Responsabilità:
  1. Creazione/caricamento presentazione
  2. Aggiunta slide con parametri, grafici, risultati
  3. Salvataggio file PPTX finale
  4. Utilities per formattazione slide
"""

import os
from pathlib import Path
from typing import List, Dict, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from config.settings import DOF_NAMES, OF_NAMES, DOF_BOUNDS, ACTIVE_DOF_INDICES
from config.paths import PRESENTATION_TEMPLATE_PATH, PRESENTATION_OUTPUT_PATH, OUTPUT_DIR
from utils.logger import logger


def load_or_create_presentation(template_path: Optional[str] = None) -> Presentation:
    """
    Carica una presentazione da template o crea una nuova.

    Args:
        template_path: Path al template PPTX. Se None o non esiste, crea nuova

    Returns:
        Presentation object
    """

    if template_path and os.path.exists(template_path):
        logger.info(f"Loading presentation template: {template_path}")
        prs = Presentation(template_path)
    else:
        logger.info("Creating new presentation from scratch")
        prs = Presentation()
        # Layout default: 16:9
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

    return prs


def add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    """Aggiunge slide titolo."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    if subtitle and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle


def add_content_slide(
        prs: Presentation,
        title: str,
        left_content: str = "",
        right_content: str = ""
) -> None:
    """
    Aggiunge slide con titolo e contenuto a sinistra/destra.

    Args:
        prs: Presentation object
        title: Titolo della slide
        left_content: Contenuto colonna sinistra (testo)
        right_content: Contenuto colonna destra (testo)
    """
    slide_layout = prs.slide_layouts[3]  # Content with caption layout
    slide = prs.slides.add_slide(slide_layout)

    if slide.shapes.title:
        slide.shapes.title.text = title

    # Colonna sinistra
    if left_content and len(slide.placeholders) > 1:
        left_box = slide.placeholders[1]
        tf = left_box.text_frame
        tf.clear()

        for line in left_content.split('\n'):
            p = tf.add_paragraph()
            p.text = line
            p.level = 0

    # Colonna destra
    if right_content and len(slide.placeholders) > 2:
        right_box = slide.placeholders[2]
        tf = right_box.text_frame
        tf.clear()

        for line in right_content.split('\n'):
            p = tf.add_paragraph()
            p.text = line
            p.level = 0


def add_image_slide(prs: Presentation, title: str, image_path: str) -> None:
    """
    Aggiunge slide con titolo e immagine.

    Args:
        prs: Presentation object
        title: Titolo della slide
        image_path: Path all'immagine
    """
    if not os.path.exists(image_path):
        logger.warning(f"Image not found: {image_path}")
        return

    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Titolo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True

    # Immagine centrata
    left = Inches(0.75)
    top = Inches(1.0)
    width = Inches(8.5)

    slide.shapes.add_picture(image_path, left, top, width=width)


def add_iteration_slide(
        prs: Presentation,
        row_idx: int,
        learning_rate: float,
        batch_size: int,
        csi_original: float,
        dof_modifications: List[str],
        image_paths: List[str]
) -> None:
    """
    Aggiunge set di slide per una iterazione di training.

    Crea:
      1. Slide parametri
      2. Slide risultati (immagine)
      3. Slide metriche actor (immagine)
      4. Slide metriche critic (immagine)

    Args:
        prs: Presentation object
        row_idx: Indice riga dataset
        learning_rate: LR usato
        batch_size: Batch size usato
        csi_original: CSI originale dataset
        dof_modifications: Lista modifiche DOF
        image_paths: [plot_results.png, plot_metrics_actor.png, plot_metrics_critic.png]
    """

    # ─── SLIDE 1: Parametri ───
    slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(slide_layout)

    if slide.shapes.title:
        slide.shapes.title.text = f"Training Iterazione — Row {row_idx} | LR={learning_rate}"

    # Sinistra: Parametri
    if len(slide.placeholders) > 1:
        left_box = slide.placeholders[1]
        tf = left_box.text_frame
        tf.clear()

        params = [
            f"Learning Rate: {learning_rate}",
            f"Batch Size: {batch_size}",
            f"CSI Originale: {csi_original:.6f}",
            "",
            "DOF Modificati:",
        ]

        for param in params:
            p = tf.add_paragraph()
            p.text = param
            p.level = 0

        for mod in dof_modifications:
            p = tf.add_paragraph()
            p.text = mod
            p.level = 1

    # Destra: Bounds
    if len(slide.placeholders) > 2:
        right_box = slide.placeholders[2]
        tf = right_box.text_frame
        tf.clear()

        p_title = tf.add_paragraph()
        p_title.text = "DOF Bounds:"
        p_title.level = 0

        for name, bounds in zip(DOF_NAMES, DOF_BOUNDS):
            p = tf.add_paragraph()
            p.text = f"{name}: [{bounds[0]:.3f}, {bounds[1]:.3f}]"
            p.level = 0
            p.font.size = Pt(9)

    # ─── SLIDE 2-4: Immagini ───
    image_titles = [
        "Risultati Training",
        "Metriche Actor (Policy)",
        "Metriche Critic (Value Function)",
    ]

    for img_path, img_title in zip(image_paths, image_titles):
        if os.path.exists(img_path):
            add_image_slide(prs, img_title, img_path)


def save_presentation(prs: Presentation, output_path: Optional[str] = None) -> str:
    """
    Salva la presentazione.

    Args:
        prs: Presentation object
        output_path: Path dove salvare. Default: config/paths.py

    Returns:
        Path file salvato
    """
    if output_path is None:
        output_path = str(PRESENTATION_OUTPUT_PATH)

    # Crea directory se necessario
    output_dir = Path(output_path).parent
    output_dir.mkdir(parents=True, exist_ok=True)

    logger.info(f"Saving presentation: {output_path}")
    prs.save(output_path)
    logger.info(f"✅ Presentation saved: {output_path}")

    return output_path


def cleanup_temp_files() -> None:
    """
    Elimina file temporanei generati durante training.

    Elimina:
      - Immagini PNG (plot_results.png, etc.)
      - CSV monitor di Gym
      - Cartelle checkpoint e log
    """
    import glob
    import shutil

    logger.info("\nCleaning temporary files...")

    # 1. Immagini
    image_files = [
        "plot_results.png",
        "plot_metrics_actor.png",
        "plot_metrics_critic.png",
    ]

    for img in image_files:
        if os.path.exists(img):
            os.remove(img)
            logger.debug(f"Removed: {img}")

    # 2. Monitor CSV
    for monitor_file in glob.glob("*monitor*.csv"):
        try:
            os.remove(monitor_file)
            logger.debug(f"Removed: {monitor_file}")
        except:
            pass

    # 3. Checkpoint e log directory
    patterns = [
        "ppo_blade_checkpoints",
        "ppo_blade_*_checkpoints_*",
        "ppo_blade_*_logs_*",
    ]

    for pattern in patterns:
        for directory in glob.glob(pattern):
            if os.path.isdir(directory):
                try:
                    shutil.rmtree(directory)
                    logger.debug(f"Removed directory: {directory}")
                except:
                    pass

    logger.info("✅ Cleanup completed!")


# ============================================================
# Utility functions
# ============================================================

def get_dof_bounds_string() -> str:
    """Restituisce stringa formattata con i bound di tutti i DOF."""
    lines = []
    for name, bounds in zip(DOF_NAMES, DOF_BOUNDS):
        lines.append(f"{name}: [{bounds[0]:.4f}, {bounds[1]:.4f}]")
    return "\n".join(lines)


def get_active_dof_string() -> str:
    """Restituisce stringa con DOF attivi."""
    active_names = [DOF_NAMES[i] for i in ACTIVE_DOF_INDICES]
    return ", ".join(active_names)


if __name__ == "__main__":
    # Test
    logger.info("Testing presentation generation...")

    prs = load_or_create_presentation()
    add_title_slide(prs, "PPO Blade Optimization", "Test Report")
    add_content_slide(
        prs,
        "Test Content",
        "Left column text",
        "Right column text"
    )

    output = save_presentation(prs)
    logger.info(f"Test presentation created: {output}")