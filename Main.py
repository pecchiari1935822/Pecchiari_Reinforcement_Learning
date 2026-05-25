import os
import pandas as pd
import numpy as np
from pptx.util import Inches, Pt
from Agente.Set_input_param import ROW_INDEX, combinazioni_da_testare, learning_rate, n_steps, ACTIVE_DOF_INDICES, OF_NAMES, TARGET_CSI
from Agente.PPO import DOF_BOUNDS_ALL, DOF_NAMES_ALL, PPO_PARAMS, SURROGATE_MODEL_PATH, SCALER_PATH, aggiungi_slide_iterazione, train, pulisci_file_temporanei, Presentation, Path
from Ambiente.Ambiente_claude_senza_keras import load_surrogate
from Smith_Chart.Reaction_total_to_total.Smith_chart_reaction_total_to_total import smith_reaction_total_to_total
from Smith_Chart.Action_total_to_static.Smith_chart_uscita_assiale import smith_action_assiale
from Smith_Chart.Action_total_to_total.Smith_chart_total_to_total import smith_action_total_to_total

# ==========================================
# TASK 1: Partenza casuale
# ==========================================
def task_1(use_delta):

    use_delta = use_delta
    episode_length = 60

    prs = Presentation()
    DATABASE_DIR = Path(__file__).parent.resolve()
    TEMPLATE_PATH = str(DATABASE_DIR / "Agente" / "Template.pptx")

    if os.path.exists(TEMPLATE_PATH):
        print(f"\nCaricamento template da: {TEMPLATE_PATH}")
        prs = Presentation(TEMPLATE_PATH)
    else:
        print(f"\n⚠️  Template non trovato. Creazione presentazione vuota.")

    results = []

    # Loop su run casuali


    # Slide intro per questo run
    slide_layout_intro = prs.slide_layouts[0]
    slide_intro = prs.slides.add_slide(slide_layout_intro)

    if slide_intro.shapes.title:
        title_shape = slide_intro.shapes.title
        title_shape.text = f"Task 1 - Run Casuale"

    if len(slide_intro.placeholders) > 1:
        subtitle_shape = slide_intro.placeholders[1]
        subtitle_shape.text = "Partenza da profilo completamente casuale"

    # Slide con parametri PPO
    slide_layout_params = prs.slide_layouts[2]
    slide_params = prs.slides.add_slide(slide_layout_params)

    if slide_params.shapes.title:
        slide_params.shapes.title.text = "Parametri PPO"

    if len(slide_params.placeholders) > 1:
        corpo_params = slide_params.placeholders[1]
        tf_params = corpo_params.text_frame
        tf_params.clear()
        for key, value in PPO_PARAMS.items():
            p = tf_params.add_paragraph()
            if key.strip() == "":
                p.text = ""
            else:
                p.text = f"{key} = {value}"

    surrogate_fn = load_surrogate(SURROGATE_MODEL_PATH, SCALER_PATH)

    # Loop su learning_rate e n_steps
    for lr in learning_rate:
        for n_step in n_steps:

            # Calcola batch_size
            if n_step <= 10:
                batch_size = 10
            elif n_step <= 50:
                batch_size = 32
            else:
                batch_size = 64

            print(f"\n  Learning_rate: {lr}, N_steps: {n_step}, Batch_size: {batch_size}")
            print(f"  {'-'*70}")

            # Esegui training
            model, best_dof, best_of, best_csi, model_path = train(
                surrogate_fn=surrogate_fn,
                start_dof=None,  # ← CASUALE
                learning_rate=lr,
                n_steps=n_step,
                batch_size=batch_size, use_delta =use_delta, episode_length=episode_length
            )

            phi_ottimale = float(best_of[OF_NAMES.index("OF_phi")])
            psi_ottimale = float(best_of[OF_NAMES.index("OF_psi")])

            smith_action_assiale.plot(target_point=(phi_ottimale, psi_ottimale), highlight_deflection=100,
                                      save_path="smith_diagram_action_assiale.png")
            smith_action_total_to_total.plot(target_point=(phi_ottimale, psi_ottimale), highlight_deflection=100,
                                             save_path="smith_diagram_action_total_to_total.png")
            smith_reaction_total_to_total.plot(target_point=(phi_ottimale, psi_ottimale), highlight_deflection=80,
                                               save_path="smith_diagram_reaction_total_to_total.png")

            print(f"  ✓ Best CSI: {best_csi:.6f}")
            print(f"  ✓ Modello: {model_path}")

            # Prepara parametri per slide
            parametri_iterazione = {
                "Tipo": "Casuale (No Dataset)",
                "Episode Length": episode_length,
                "N steps": n_step,
                "Batch size": batch_size,
                "Learning Rate": lr,
                "Best CSI trovato": round(best_csi, 6),
                "DOF attivi modificati": [f"{DOF_NAMES_ALL[i]}" for i in ACTIVE_DOF_INDICES]
            }

            # Immagini generate dal training
            img_paths = [
                "plot_results.png",
                "plot_dof_evolution.png",
                "plot_dof_evolution_barre.png",
                "plot_metrics_actor.png",
                "plot_metrics_critic.png",

            ]

            # Aggiungi slide con risultati
            # NOTA: start_dof=None perché non sappiamo il profilo iniziale casuale
            try:
                aggiungi_slide_iterazione(
                    prs,
                    parametri_iterazione,
                    img_paths,
                    row_idx=None,  # Non c'è una riga specifica del dataset
                    lr=lr,
                    best_dof=best_dof,
                    best_of=best_of,
                    start_dof=None,  # Placeholder
                    start_of=None
                )
            except Exception as e:
                print(f"  ⚠️  Errore nell'aggiunta slide: {e}")

            slide = prs.slides.add_slide(prs.slide_layouts[4])
            if slide.shapes.title:
                slide.shapes.title.text = "Smith Diagram Action - Axial exit"
            slide.shapes.add_picture("smith_diagram_action_assiale.png", Inches(1.5), Inches(1.2), height=Inches(5.2))

            # Slide 2: Smith Diagram - Action Total to Total
            slide = prs.slides.add_slide(prs.slide_layouts[4])
            if slide.shapes.title:
                slide.shapes.title.text = "Smith Diagram Action - Total to Total"
            slide.shapes.add_picture("smith_diagram_action_total_to_total.png", Inches(1.5), Inches(1.2), height=Inches(5.2))

            # Slide 3: Smith Diagram - Reaction Total to Total
            slide = prs.slides.add_slide(prs.slide_layouts[4])
            if slide.shapes.title:
                slide.shapes.title.text = "Smith Diagram Reaction - Total to Total"
            slide.shapes.add_picture("smith_diagram_reaction_total_to_total.png", Inches(1.5), Inches(1.2), height=Inches(5.2))

            # Salva risultati
            results.append({
                'learning_rate': lr,
                'n_steps': n_step,
                'batch_size': batch_size,
                'best_csi': best_csi,
                'model_path': model_path
            })



    # Salva presentazione
    if use_delta == True:
        output_pptx = "Task1_delta.pptx"
    else:
        output_pptx = "Task1_mapping completo.pptx"
    prs.save(output_pptx)

    # Salva risultati in CSV
    df_results = pd.DataFrame(results)
    csv_path = "task1_results.csv"
    df_results.to_csv(csv_path, index=False)

    print("\n" + "="*70)
    print("  TASK 1 COMPLETATO!")
    print("="*70)
    print(f"\n  📊 Presentazione salvata: {output_pptx}")

    print(f"\n  Riepilogo:")

    print(f"    - Miglior CSI globale: {df_results['best_csi'].min():.6f}")
    print(f"    - CSI medio: {df_results['best_csi'].mean():.6f}")
    print(f"    - Std Dev: {df_results['best_csi'].std():.6f}")
    print("\n")

    pulisci_file_temporanei()


# ==========================================
# TASK 2: Ottimizzazione da riga del Dataset
# ==========================================
def task_2(use_delta):

    use_delta = use_delta
    if use_delta == True:
        print("\n⚠️  Modalità DELTA attiva: il PPO ottimizzerà la differenza rispetto al CSI originale.")
    else:
        print ("\n⚠️  Modalità mapping completo: il PPO ottimizzerà direttamente il CSI senza considerare il delta.")
    episode_length = 20

    prs = Presentation()

    # 1. Imposta il percorso del tuo dataset e la riga che vuoi analizzare
    DATABASE_DIR = Path(__file__).parent.resolve()
    DATASET_PATH = str(DATABASE_DIR / "Data" / "database.dat")
    TEMPLATE_PATH = str(DATABASE_DIR / "Agente" / "Template.pptx")

    if os.path.exists(TEMPLATE_PATH):
        print(f"Caricamento template da: {TEMPLATE_PATH}")
        prs = Presentation(TEMPLATE_PATH)

    df = pd.read_csv(DATASET_PATH)
    df.columns = df.columns.str.replace("_OP_01", "", regex=False)
    df.columns = df.columns.str.replace("_GEOM_", "", regex=False)
    df.columns = df.columns.str.replace("_BC_", "", regex=False)

    for row_idx in ROW_INDEX:
        print(f"\nLettura dataset: {DATASET_PATH}")
        print(f"Estrazione riga numero: {row_idx}")

        # Estrai la riga dal dataset

        row = df.iloc[row_idx]
        riga = df.iloc[row_idx].values

        # 3. Estrai il CSI originale (12° valore -> indice 11)
        csi_originale = float(row[TARGET_CSI])

        for active_dof in combinazioni_da_testare:
            # ---> FIX CRITICO: Aggiorniamo dinamicamente le variabili dell'ambiente!
            # Altrimenti Gym continuerebbe ad usare i parametri fissi del file .py
            import Ambiente.Ambiente_claude_senza_keras as env_module

            env_module.ACTIVE_DOF_INDICES = active_dof
            env_module.DOF_BOUNDS = [DOF_BOUNDS_ALL[i] for i in active_dof]

            # Aggiorniamo anche la variabile globale in questo script per i grafici

            ACTIVE_DOF_INDICES = active_dof

            slide_layout_intro = prs.slide_layouts[0]
            slide_intro = prs.slides.add_slide(slide_layout_intro)

            if slide_intro.shapes.title:
                # Prendi direttamente la forma (shape) del titolo
                title_shape = slide_intro.shapes.title
                title_shape.text = f"Ottimizzazione Riga Dataset: {row_idx}"

            # --- 2. GESTIONE DEL SOTTOTITOLO ---
            # Verifica che esista un secondo placeholder (generalmente usato per il sottotitolo)
            if len(slide_intro.placeholders) > 1:
                subtitle_shape = slide_intro.placeholders[1]

                nomi_attivi = [DOF_NAMES_ALL[i].replace("DOF_", "").replace("_GEOM", "").replace("_", "") for i in active_dof]
                subtitle_shape.text = f"Analisi {' e '.join(nomi_attivi).lower()}"

                # 0. Slide per il proiflo alla riga n del dataset
                slide_layout_testo_ppo = prs.slide_layouts[2]
                slide_testo_ppo = prs.slides.add_slide(slide_layout_testo_ppo)
                slide_testo_ppo.shapes.title.text = f"Parametri PPO"
                if len(slide_testo_ppo.placeholders) > 1:
                    corpo_ppo = slide_testo_ppo.placeholders[1]
                    tf_ppo = corpo_ppo.text_frame
                    tf_ppo.clear()  # Pulisce il testo di default
                    for key, value in PPO_PARAMS.items():
                        p = tf_ppo.add_paragraph()
                        run = p.add_run()

                        if key.strip() == "":
                            run.text = ""  # riga vuota
                        else:
                            run.text = f"{key} = {value}"

            # 4. Poni a 0.0 i DOF attivi (quelli che il PPO deve modificare)
            print("\n--- PERTURBAZIONE DEI DOF ATTIVI ---")
            print(f"Estrazione riga numero: {row_idx}")

            # 2. Estrai i 7 DOF (dal 3° al 9° valore -> indici da 2 a 8)
            start_profile = row[DOF_NAMES_ALL].values.astype(np.float32).copy()
            start_of_originali = row[OF_NAMES].values.astype(np.float32)

            # Perturbazione del dof attivo (che si vuole ottimizzare)
            for idx in active_dof:
                val_originale = start_profile[idx]
                min_bound, max_bound = DOF_BOUNDS_ALL[idx]

                # Calcoliamo la distanza tra il valore originale e i due limiti
                distanza_dal_min = abs(val_originale - min_bound)
                distanza_dal_max = abs(max_bound - val_originale)

                # Scegliamo il limite più lontano per rendere la sfida più difficile per il PPO!
                if distanza_dal_max > distanza_dal_min:
                    valore_perturbato = max_bound
                else:
                    valore_perturbato = min_bound

                # Applichiamo la perturbazione
                start_profile[idx] = valore_perturbato

                surrogate = load_surrogate(SURROGATE_MODEL_PATH, SCALER_PATH)

                input_dof = start_profile.reshape(1, -1)
                predizione = surrogate(input_dof)
                csi_modificato = float(predizione[11])

                nome_dof = DOF_NAMES_ALL[idx]
                print(f"  {nome_dof:<28}: {val_originale:>10.6f}  -->  Spostato a: {valore_perturbato:>10.6f}")
                print(f"  (CSI originale: {csi_originale:.6f}  -->  Perturbato a: {csi_modificato:.6f})")

            print("-------------------------------------------------")
            print("\nI DOF attivi sono stati spostati ai loro limiti per testare il PPO.")

            print(f"\nProfilo di partenza per il PPO: {start_profile}.3f")

            surrogate_fn = load_surrogate(SURROGATE_MODEL_PATH, SCALER_PATH)

            # Testo l'addestramento di ogni profilo (riga) con diversi learning_rate per vedere quale è meglio
            for lr in learning_rate:

                # Per ogni learning_rate, provo anche diversi n_steps per vedere l'effetto sulla convergenza e sul risultato finale
                for n_step in n_steps:
                    if n_step <= 10:
                        batch_size = 10
                    elif n_step <= 50:
                        batch_size = 32
                    else:
                        batch_size = 64

                    # 5. Avvia il Training (Task 2)
                    print(f"\n  Learning_rate attuale : {lr}\n")
                    model, best_dof, best_of, best_csi, model_ = train(
                        surrogate_fn=surrogate_fn,
                        start_dof=start_profile,
                        learning_rate=lr, n_steps=n_step, batch_size=batch_size, ROW_INDEX=ROW_INDEX, use_delta =use_delta, episode_length=episode_length,
                        ref_of=start_of_originali
                    )

                    phi_ottimale = float(best_of[OF_NAMES.index("OF_phi")])
                    psi_ottimale = float(best_of[OF_NAMES.index("OF_psi")])

                    smith_action_assiale.plot(target_point=(phi_ottimale, psi_ottimale), highlight_deflection=100,
                                              save_path="smith_diagram_action_assiale.png")
                    smith_action_total_to_total.plot(target_point=(phi_ottimale, psi_ottimale),
                                                     highlight_deflection=100,
                                                     save_path="smith_diagram_action_total_to_total.png")
                    smith_reaction_total_to_total.plot(target_point=(phi_ottimale, psi_ottimale),
                                                       highlight_deflection=80,
                                                       save_path="smith_diagram_reaction_total_to_total.png")




                    # 6. Confronto Finale PPO vs Dataset
                    miglioramento = csi_originale - best_csi
                    segno = "✓ MIGLIORATO" if miglioramento > 0 else "✗ PEGGIORATO"

                    print("\n" + "=" * 60)
                    print("  RISULTATO TASK 2: CONFRONTO PPO vs DATASET")
                    print("=" * 60)
                    print(f"  CSI Originale (dataset) : {csi_originale:.6f}")
                    print(f"  CSI Migliore PPO        : {best_csi:.6f}")
                    print(f"  Differenza (Delta CSI)  : {miglioramento:+.6f}  ({segno})")
                    print("=" * 60)

                    print(f"\n  Riga {row_idx} (index={row.name}):")
                    print(row.to_frame().T.to_string(index=False))

                    path_img1 = "plot_results.png"
                    path_img2 = "plot_dof_evolution.png"
                    path_img3 = "plot_dof_evolution_barre.png"
                    path_img4 = "plot_metrics_actor.png"
                    path_img5 = "plot_metrics_critic.png"


                    parametri_iterazione = {
                        "N steps": n_step,
                        "Episode Length": episode_length,
                        "Batch size": batch_size,
                        "CSI Originale": round(csi_originale, 6),
                        "DOF attivi modificati": [f"{DOF_NAMES_ALL[i]} da {row[DOF_NAMES_ALL[i]]} a {start_profile[i]:.3f}" for i in active_dof],
                        "CSI profilo modificato": None
                    }


                    # 3. Aggiungi i risultati alla presentazione
                    img_paths = [path_img1, path_img2, path_img3, path_img4, path_img5]



                    aggiungi_slide_iterazione(prs, parametri_iterazione, img_paths, row_idx, lr, best_dof, best_of, start_profile, start_of_originali)

                    slide = prs.slides.add_slide(prs.slide_layouts[4])
                    if slide.shapes.title:
                        slide.shapes.title.text = "Smith Diagram Action - Axial exit"
                    slide.shapes.add_picture("smith_diagram_action_assiale.png", Inches(1.5), Inches(1.2),
                                             height=Inches(5.2))

                    # Slide 2: Smith Diagram - Action Total to Total
                    slide = prs.slides.add_slide(prs.slide_layouts[4])
                    if slide.shapes.title:
                        slide.shapes.title.text = "Smith Diagram Action - Total to Total"
                    slide.shapes.add_picture("smith_diagram_action_total_to_total.png", Inches(1.5), Inches(1.2),
                                             height=Inches(5.2))

                    # Slide 3: Smith Diagram - Reaction Total to Total
                    slide = prs.slides.add_slide(prs.slide_layouts[4])
                    if slide.shapes.title:
                        slide.shapes.title.text = "Smith Diagram Reaction - Total to Total"
                    slide.shapes.add_picture("smith_diagram_reaction_total_to_total.png", Inches(1.5), Inches(1.2),
                                             height=Inches(5.2))

        if use_delta == True:
            output_pptx = f"Task2_delta_riga{row_idx}.pptx"
        else:
            output_pptx = f"Task2_mapping_completo_riga{row_idx}.pptx"
        prs.save(output_pptx)

    pulisci_file_temporanei()

if __name__ == "__main__":
    #task_1(True)
    #task_1(False)
    task_2(True)
    #task_2(False)