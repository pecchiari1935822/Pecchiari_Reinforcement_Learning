import os
import pandas as pd
import numpy as np
from pptx import Presentation
from pathlib import Path
from Config.Set_input_param import ROW_INDEX, combinazioni_da_testare, learning_rate, n_steps, ACTIVE_DOF_INDICES, \
    OF_NAMES, TARGET_CSI, perturbazione_dof_attivi, DOF_BOUNDS_ALL, DOF_NAMES_ALL, dataset, df, USE_MULTIMODEL
from Agente.FOCOPS import train, pulisci_file_temporanei
from Ambiente.FOCOPS_Ambiente import  surrogate
from Report.Presentazione import aggiungi_slide_iterazione, slide_iniziali_task_1, slide_iniziali_task_2, aggiungi_smith
from Report.Plot import plot_smith, plot_smith_zoom_reaction

# ==========================================
# TASK 1: Partenza casuale
# ==========================================
def task_1(use_delta):

    use_delta = use_delta
    episode_length = 60
    task_1 = True

    prs = Presentation()
    DATABASE_DIR = Path(__file__).parent.resolve()
    TEMPLATE_PATH = str(DATABASE_DIR / "Report" / "Template.pptx")

    if os.path.exists(TEMPLATE_PATH):
        print(f"\nCaricamento template da: {TEMPLATE_PATH}")
        prs = Presentation(TEMPLATE_PATH)
    else:
        print(f"\n⚠️  Template non trovato. Creazione presentazione vuota.")

    results = []

    # Slide intro per questo run
    try:
        slide_iniziali_task_1(prs)
    except Exception as e:
        print(f"  ⚠️  Errore nell'aggiunta slide iniali: {e}")




    # Loop su learning_rate e n_steps
    for lr in learning_rate:
        for n_step in n_steps:

            # Calcola batch_size
            if n_step <= 10:
                batch_size = 10
            elif n_step <= 50:
                batch_size = 32
            else:
                batch_size = 64

            print(f"\n  Learning_rate: {lr}, N_steps: {n_step}, Batch_size: {batch_size}")
            print(f"  {'-'*70}")

            # Esegui training
            model, best_dof, best_of, best_csi, model_path = train(

                start_dof=None,  # ← CASUALE
                learning_rate=lr,
                n_steps=n_step,
                batch_size=batch_size, use_delta =use_delta, episode_length=episode_length, task1=task_1
            )

            if best_of is None:
                raise RuntimeError(
                    "best_of è None: nessun profilo è stato salvato come best. Controlla callback/env info['of'].")

            # Cerca l'indice di PHI e PSI in modo flessibile nella lista OF_NAMES
            idx_phi_reale = next((i for i, name in enumerate(OF_NAMES) if "PHI" in name.upper()), None)
            idx_psi_reale = next((i for i, name in enumerate(OF_NAMES) if "PSI" in name.upper()), None)
            idx_alfa_ex_reale = next((i for i, name in enumerate(OF_NAMES) if "ALFA_EX" in name.upper()), None)
            idx_alfa_in_reale = next((i for i, name in enumerate(DOF_NAMES_ALL) if "ALFAIN" in name.upper()), None)
            idx_beta1_reale = next((i for i, name in enumerate(DOF_NAMES_ALL) if "BETA1" in name.upper()), None)
            idx_beta2_reale = next((i for i, name in enumerate(DOF_NAMES_ALL) if "BETA2" in name.upper()), None)

            if idx_phi_reale is not None:
                phi_ottimale = float(best_of[idx_phi_reale])

            if idx_psi_reale is not None:
                psi_ottimale = float(best_of[idx_psi_reale])


            alpha_ex_ottimale = float(best_of[idx_alfa_ex_reale])
            alpha_in_ottimale = float(best_dof[idx_alfa_in_reale])
            beta_1_ottimale = float(best_dof[idx_beta1_reale])
            beta_2_ottimale = float(best_dof[idx_beta2_reale])



            plot_smith(
                phi_ottimale=phi_ottimale,
                psi_ottimale=psi_ottimale,
                defl_min=40,
                defl_max=140,
                step=1,  # ← Interpola ogni 1°
                validate=True
            )

            print(f"  ✓ Best CSI: {best_csi:.6f}")
            print(f"  ✓ Modello: {model_path}")

            # Prepara parametri per slide
            parametri_iterazione = {
                "Tipo": "Casuale (No Dataset)",
                "Episode Length": episode_length,
                "N steps": n_step,
                "Batch size": batch_size,
                "Learning Rate": lr,
                "Best CSI trovato": round(best_csi, 6),
                "DOF attivi modificati": [f"{DOF_NAMES_ALL[i]}" for i in ACTIVE_DOF_INDICES]
            }

            # Immagini generate dal training
            import glob
            img_paths = ["plot_results.png"]
            img_paths.extend(sorted(glob.glob("plot_dof_evolution_*.png")))
            img_paths.extend([
                "plot_metrics_actor.png",
                "plot_metrics_critic.png",
            ])

            # Aggiungi slide con risultati
            # NOTA: start_dof=None perché non sappiamo il profilo iniziale casuale
            try:
                aggiungi_slide_iterazione(
                    prs,
                    parametri_iterazione,
                    img_paths,
                    row_idx=None,  # Non c'è una riga specifica del dataset
                    lr=lr,
                    best_dof=best_dof,
                    best_of=best_of,
                    start_dof=None,  # Placeholder
                    start_of=None,
                    alpha_ex_ottimale=alpha_ex_ottimale,
                    alpha_in_ottimale=alpha_in_ottimale,
                    beta_1_ottimale=beta_1_ottimale,
                    beta_2_ottimale=beta_2_ottimale
                )
            except Exception as e:
                print(f"  ⚠️  Errore nell'aggiunta slide: {e}")

            aggiungi_smith(prs, task1=task_1)
            # Salva risultati
            results.append({
                'learning_rate': lr,
                'n_steps': n_step,
                'batch_size': batch_size,
                'best_csi': best_csi,
                'model_path': model_path
            })



    # Salva presentazione
    if use_delta == True:
        output_pptx = os.path.join("Risultati", "Presentazioni", "Task1_delta.pptx")
    else:
        output_pptx = os.path.join("Risultati", "Presentazioni", "Task1_mapping completo.pptx")
    prs.save(output_pptx)

    # Salva risultati in CSV
    df_results = pd.DataFrame(results)
    csv_path = "task1_results.csv"
    df_results.to_csv(csv_path, index=False)

    pulisci_file_temporanei(task_1=task_1)


# ==========================================
# TASK 2: Ottimizzazione da riga del Dataset
# ==========================================
def task_2(use_delta):

    use_delta = use_delta
    if use_delta == True:
        print("\n⚠️  Modalità DELTA attiva: il PPO ottimizzerà la differenza rispetto al CSI originale.")
    else:
        print ("\n⚠️  Modalità mapping completo: il PPO ottimizzerà direttamente il CSI senza considerare il delta.")
    episode_length = 20

    # 1. Imposta il percorso del tuo dataset e la riga che vuoi analizzare
    DATABASE_DIR = Path(__file__).parent.resolve()
    DATASET_PATH = dataset
    TEMPLATE_PATH = str(DATABASE_DIR / "Report" / "Template.pptx")

    if os.path.exists(TEMPLATE_PATH):
        print(f"Caricamento template da: {TEMPLATE_PATH}")
        prs = Presentation(TEMPLATE_PATH)



    for row_idx in ROW_INDEX:

        if os.path.exists(TEMPLATE_PATH):
            print(f"Caricamento template da: {TEMPLATE_PATH}")
            prs = Presentation(TEMPLATE_PATH)
        else:
            prs = Presentation()

        print(f"\nLettura dataset: {DATASET_PATH}")
        print(f"Estrazione riga numero: {row_idx}")

        # Estrai la riga dal dataset

        row = df.iloc[row_idx]
        riga = df.iloc[row_idx].values

        # 3. Estrai il CSI originale (12° valore -> indice 11)
        csi_originale = float(row[TARGET_CSI])

        for active_dof in combinazioni_da_testare:
            # ---> FIX CRITICO: Aggiorniamo dinamicamente le variabili dell'ambiente!
            # Altrimenti Gym continuerebbe ad usare i parametri fissi del file .py
            import Ambiente.PPO_Ambiente as env_module

            env_module.ACTIVE_DOF_INDICES = active_dof
            env_module.DOF_BOUNDS = [DOF_BOUNDS_ALL[i] for i in active_dof]

            # Aggiorniamo anche la variabile globale in questo script per i grafici

            ACTIVE_DOF_INDICES = active_dof

            try:
                slide_iniziali_task_2(prs, row_idx, active_dof)
            except Exception as e:
                print(f"  ⚠️  Errore nell'aggiunta slide iniali: {e}")

            print(f"Estrazione riga numero: {row_idx}")

            # 2. Estrai i 7 DOF (dal 3° al 9° valore -> indici da 2 a 8)
            start_profile = row[DOF_NAMES_ALL].values.astype(np.float32).copy()
            start_of_originali = row[OF_NAMES].values.astype(np.float32)
            start_dof_dataset = start_profile.copy()

            # Perturbazione del dof attivo (che si vuole ottimizzare)
            if perturbazione_dof_attivi == True:
                print("\n--- PERTURBAZIONE DEI DOF ATTIVI ---")
                for idx in active_dof:
                    val_originale = start_profile[idx]
                    min_bound, max_bound = DOF_BOUNDS_ALL[idx]

                    # Calcoliamo la distanza tra il valore originale e i due limiti
                    distanza_dal_min = abs(val_originale - min_bound)
                    distanza_dal_max = abs(max_bound - val_originale)

                    # Scegliamo il limite più lontano per rendere la sfida più difficile per il PPO!
                    if distanza_dal_max > distanza_dal_min:
                        valore_perturbato = max_bound
                    else:
                        valore_perturbato = min_bound

                    # Applichiamo la perturbazione
                    start_profile[idx] = valore_perturbato



                    input_dof = start_profile.reshape(1, -1)
                    predizione = surrogate(input_dof)
                    csi_modificato = float(predizione[11])

                    nome_dof = DOF_NAMES_ALL[idx]
                    print(f"  {nome_dof:<28}: {val_originale:>10.6f}  -->  Spostato a: {valore_perturbato:>10.6f}")
                    print(f"  (CSI originale: {csi_originale:.6f}  -->  Perturbato a: {csi_modificato:.6f})")

                print("-------------------------------------------------")
                print("\nI DOF attivi sono stati spostati ai loro limiti per testare il PPO.")

                print(f"\nProfilo di partenza per il PPO: {start_profile}.3f")



            # Testo l'addestramento di ogni profilo (riga) con diversi learning_rate per vedere quale è meglio
            for lr in learning_rate:

                # Per ogni learning_rate, provo anche diversi n_steps per vedere l'effetto sulla convergenza e sul risultato finale
                for n_step in n_steps:
                    if n_step <= 10:
                        batch_size = 10
                    elif n_step <= 50:
                        batch_size = 32
                    else:
                        batch_size = 64

                    # 5. Avvia il Training (Task 2)
                    print(f"\n  Learning_rate attuale : {lr}\n")
                    model, best_dof, best_of, best_csi, model_ = train(

                        start_dof=start_profile,
                        learning_rate=lr, n_steps=n_step, batch_size=batch_size, ROW_INDEX=ROW_INDEX, use_delta =use_delta, episode_length=episode_length,
                        ref_of=start_of_originali
                    )

                    print(best_of)
                    if best_of is None:
                        raise RuntimeError(
                            "best_of è None: nessun profilo è stato salvato come best. Controlla callback/env info['of'].")

                    # Cerca l'indice di PHI e PSI in modo flessibile nella lista OF_NAMES
                    idx_phi_reale = next((i for i, name in enumerate(OF_NAMES) if "PHI" in name.upper()), None)
                    idx_psi_reale = next((i for i, name in enumerate(OF_NAMES) if "PSI" in name.upper()), None)
                    idx_alfa_ex_reale = next((i for i, name in enumerate(OF_NAMES) if "ALFA_EX" in name.upper()), None)
                    idx_alfa_in_reale = next((i for i, name in enumerate(DOF_NAMES_ALL) if "ALFAIN" in name.upper()), None)
                    idx_beta1_reale = next((i for i, name in enumerate(DOF_NAMES_ALL) if "BETA1" in name.upper()), None)
                    idx_beta2_reale = next((i for i, name in enumerate(DOF_NAMES_ALL) if "BETA2" in name.upper()), None)

                    if idx_phi_reale is not None:
                        phi_ottimale = float(best_of[idx_phi_reale])

                    if idx_psi_reale is not None:
                        psi_ottimale = float(best_of[idx_psi_reale])

                    alpha_ex_ottimale = float(best_of[idx_alfa_ex_reale])

                    alpha_in_ottimale = float(best_dof[idx_alfa_in_reale])

                    beta_1_ottimale = float(best_dof[idx_beta1_reale])
                    beta_2_ottimale = float(best_dof[idx_beta2_reale])

                    orig_phi = float(start_of_originali[idx_phi_reale])
                    orig_psi = float(start_of_originali[idx_psi_reale])



                    plot_smith(
                        phi_ottimale=phi_ottimale,
                        psi_ottimale=psi_ottimale,
                        defl_min=40,
                        defl_max=140,
                        step=1,  # ← Interpola ogni 1°
                        validate=True
                    )

                    plot_smith_zoom_reaction(phi_ottimale, psi_ottimale, orig_phi, orig_psi)

                    # 6. Confronto Finale PPO vs Dataset
                    miglioramento = csi_originale - best_csi
                    segno = "✓ MIGLIORATO" if miglioramento > 0 else "✗ PEGGIORATO"

                    print("\n" + "=" * 60)
                    print("  RISULTATO TASK 2: CONFRONTO PPO vs DATASET")
                    print("=" * 60)
                    print(f"  CSI Originale (dataset) : {csi_originale:.6f}")
                    print(f"  CSI Migliore PPO        : {best_csi:.6f}")
                    print(f"  Differenza (Delta CSI)  : {miglioramento:+.6f}  ({segno})")
                    print("=" * 60)

                    print(f"\n  Riga {row_idx} (index={row.name}):")
                    print(row.to_frame().T.to_string(index=False))

                    import glob

                    path_img1 = "plot_results.png"
                    # glob.glob espande l'asterisco e trova plot_dof_evolution_1.png, _2.png, _3.png
                    path_dof_evolutions = sorted(glob.glob("plot_dof_evolution_*.png"))
                    path_img4 = "plot_metrics_actor.png"
                    path_img5 = "plot_metrics_critic.png"

                    # Costruiamo la lista unendo dinamicamente i grafici trovati
                    img_paths = [path_img1] + path_dof_evolutions + [path_img4, path_img5]


                    parametri_iterazione = {
                        "N steps": n_step,
                        "Episode Length": episode_length,
                        "Batch size": batch_size,
                        "CSI Originale": round(csi_originale, 6),
                        "DOF attivi modificati": [f"{DOF_NAMES_ALL[i]} da {row[DOF_NAMES_ALL[i]]} a {start_profile[i]:.3f}" for i in active_dof],
                        "CSI profilo modificato": None
                    }


                    aggiungi_slide_iterazione(prs, parametri_iterazione, img_paths, row_idx, lr, best_dof, best_of, start_dof_dataset, start_of_originali,
                                              alpha_ex_ottimale, alpha_in_ottimale, beta_1_ottimale, beta_2_ottimale)

                    aggiungi_smith(prs, task1=None)

        if use_delta == True:
            output_pptx = os.path.join("Risultati", "Presentazioni", f"Task2_delta_riga{row_idx}.pptx")
        else:
            output_pptx = os.path.join("Risultati", "Presentazioni", f"Task2_mapping_completo_riga{row_idx}.pptx")
        prs.save(output_pptx)

    pulisci_file_temporanei()

if __name__ == "__main__":
    #task_1(True)
    #task_1(False)
    task_2(True)
    #task_2(False)