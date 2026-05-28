from Config.Set_input_param import DOF_NAMES_ALL, DOF_BOUNDS_ALL, OF_NAMES, ACTIVE_DOF_INDICES
from pptx.util import Inches
import os
from Config.Set_input_param import PPO_PARAMS



def slide_iniziali_task_1(prs):
    slide_layout_intro = prs.slide_layouts[0]
    slide_intro = prs.slides.add_slide(slide_layout_intro)

    if slide_intro.shapes.title:
        title_shape = slide_intro.shapes.title
        title_shape.text = f"Task 1 - Run Casuale"

    if len(slide_intro.placeholders) > 1:
        subtitle_shape = slide_intro.placeholders[1]
        subtitle_shape.text = "Partenza da profilo completamente casuale"

    # Slide con parametri PPO
    slide_layout_params = prs.slide_layouts[2]
    slide_params = prs.slides.add_slide(slide_layout_params)

    if slide_params.shapes.title:
        slide_params.shapes.title.text = "Parametri PPO"

    if len(slide_params.placeholders) > 1:
        corpo_params = slide_params.placeholders[1]
        tf_params = corpo_params.text_frame
        tf_params.clear()
        for key, value in PPO_PARAMS.items():
            p = tf_params.add_paragraph()
            if key.strip() == "":
                p.text = ""
            else:
                p.text = f"{key} = {value}"


def slide_iniziali_task_2(prs, row_idx, active_dof):
    slide_layout_intro = prs.slide_layouts[0]
    slide_intro = prs.slides.add_slide(slide_layout_intro)

    if slide_intro.shapes.title:
        # Prendi direttamente la forma (shape) del titolo
        title_shape = slide_intro.shapes.title
        title_shape.text = f"Ottimizzazione Riga Dataset: {row_idx}"

    # --- 2. GESTIONE DEL SOTTOTITOLO ---
    # Verifica che esista un secondo placeholder (generalmente usato per il sottotitolo)
    if len(slide_intro.placeholders) > 1:
        subtitle_shape = slide_intro.placeholders[1]

        nomi_attivi = [DOF_NAMES_ALL[i].replace("DOF_", "").replace("_GEOM", "").replace("_", "") for i in active_dof]
        subtitle_shape.text = f"Analisi {' e '.join(nomi_attivi).lower()}"

        # 0. Slide per il proiflo alla riga n del dataset
        slide_layout_testo_ppo = prs.slide_layouts[2]
        slide_testo_ppo = prs.slides.add_slide(slide_layout_testo_ppo)
        slide_testo_ppo.shapes.title.text = f"Parametri PPO"
        if len(slide_testo_ppo.placeholders) > 1:
            corpo_ppo = slide_testo_ppo.placeholders[1]
            tf_ppo = corpo_ppo.text_frame
            tf_ppo.clear()  # Pulisce il testo di default
            for key, value in PPO_PARAMS.items():
                p = tf_ppo.add_paragraph()
                run = p.add_run()

                if key.strip() == "":
                    run.text = ""  # riga vuota
                else:
                    run.text = f"{key} = {value}"

def aggiungi_slide_iterazione(prs, parametri, img_paths, row_idx, lr, best_dof, best_of,start_dof, start_of):
    """
    Aggiunge 4 slide alla presentazione: 1 di testo e 3 di immagini.
    """
    slide_layout_testo = prs.slide_layouts[3]
    slide_testo = prs.slides.add_slide(slide_layout_testo)

    if slide_testo.shapes.title:
        title_shape = slide_testo.shapes.title
        title_shape.text = f"Punto di partenza simulazione con LR = {lr}"


    # 1. Slide di testo con i parametri

    if len(slide_testo.placeholders) > 1:
        corpo_sinistro = slide_testo.placeholders[1]
        tf_sinistra = corpo_sinistro.text_frame
        tf_sinistra.clear()  # Pulisce il testo di default
        for key, value in parametri.items():
            if isinstance(value, list):
                # 1. Aggiungi la voce principale (es. "DOF attivi modificati:")
                p_main = tf_sinistra.add_paragraph()
                p_main.text = f"{key}:"
                p_main.level = 0

                # 2. Aggiungi i vari elementi della lista indentati (level = 1)
                for item in value:
                    p_sub = tf_sinistra.add_paragraph()
                    p_sub.text = str(item)
                    p_sub.level = 1  # QUESTO CREA L'INDENTAZIONE!

            # Se è un valore normale (testo o numero), lo mette come punto principale
            else:
                p = tf_sinistra.add_paragraph()
                p.level = 0

                if key.strip() == "":
                    p.text = ""  # riga vuota
                else:
                    p.text = f"{key} = {value}"



        # COLONNA DESTRA (Limiti dei DOF)
        # =======================================================
        # Posizioniamo una nuova casella di testo nella metà destra (da 7 pollici in poi)
        if len(slide_testo.placeholders) > 2:
            corpo_destro = slide_testo.placeholders[2]
            tf_destra = corpo_destro.text_frame
            tf_destra.clear()  # Cancella la scritta "Fare clic per inserire testo"

            p_title_destra = tf_destra.add_paragraph()
            p_title_destra.text = "Limiti operativi (DOF Bounds):"
            p_title_destra.level = 0

            # Aggiungiamo l'elenco dei limiti
            for nome_dof, bounds in zip(DOF_NAMES_ALL, DOF_BOUNDS_ALL):
                p_bound = tf_destra.add_paragraph()
                p_bound.text = f"{nome_dof} = [{bounds[0]}, {bounds[1]}]"


    # 2. Slide per ogni immagine (3 immagini)
    try:
        slide_layout_img = prs.slide_layouts[4]  # Prova con Solo Titolo
    except:
        slide_layout_img = prs.slide_layouts[1]  # Fallback

    titoli_immagini = ["Risultati","Andamento Valore Miglior DOF per episodio", "Andamento Valore Miglior DOF per episodio", "Metriche Attore", "Metriche Critico"]  # Layout Solo Titolo o Vuota
    for idx, img_path in enumerate(img_paths):
        if not os.path.exists(img_path):
            continue

        slide_img = prs.slides.add_slide(slide_layout_img)

        if slide_img.shapes.title:
            slide_img.shapes.title.text = f"{titoli_immagini[idx]}"

        # Inserisci immagine centrata (adattata per 16:9)
        # Una slide 16:9 tipica è larga 13.33 pollici e alta 7.5 pollici
        if idx == 0 or idx == 4:
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(11.33)  # Lascia 1 pollice di margine per lato
            slide_img.shapes.add_picture(str(img_path), left, top, width=width)
        else:
            left = Inches(1)
            top = Inches(1.1)
            heigh = Inches(5.9)  # Lascia 1 pollice di margine per lato
            slide_img.shapes.add_picture(str(img_path), left, top, height=heigh)



    slide_layout_best = prs.slide_layouts[5]
    slide_best = prs.slides.add_slide(slide_layout_best)

    if slide_best.shapes.title:
        title_shape = slide_best.shapes.title
        title_shape.text = f"Miglior Profilo Ottimizzato"


    # COLONNA SINISTRA: I 7 DOF ottimali
    if len(slide_best.placeholders) > 1:
        corpo_sin_best = slide_best.placeholders[1]
        tf_sin_best = corpo_sin_best.text_frame
        tf_sin_best.clear()

        p_title_dof = tf_sin_best.add_paragraph()
        p_title_dof.text = "DOF Ottimali:"
        p_title_dof.level = 0

        for i, (nome, val) in enumerate(zip(DOF_NAMES_ALL, best_dof)):
            p_dof = tf_sin_best.add_paragraph()

            # Se il DOF era attivo, mostra anche il valore di partenza
            if i in ACTIVE_DOF_INDICES:
                if start_dof is not None:
                    p_dof.text = f"{nome} (*) = {val:.6f} da {start_dof[i]:.6f}"
                else:
                    p_dof.text = f"{nome} (*) = {val:.6f}"
            else:
                p_dof.text = f"{nome} = {val:.6f}"
            p_dof.level = 0

    # COLONNA DESTRA: I 15 OF corrispondenti
    if len(slide_best.placeholders) > 2:
        corpo_des_best = slide_best.placeholders[2]
        tf_des_best = corpo_des_best.text_frame
        tf_des_best.clear()

        p_title_of = tf_des_best.paragraphs[0]
        p_title_of.text = "OF Risultanti:"
        p_title_of.level = 0

        try:
            # === RECUPERA I VALORI PER I CALCOLI ===
            # Esempio: OF_alfa_ex_OP_01 è il nome usato nel DB per l'angolo di uscita, lo recupero in best_of
            idx_alpha_ex = OF_NAMES.index("OF_alfa_ex")
            alpha_ex_deg = best_of[idx_alpha_ex]

            # alpha0 fisso a 10 gradi come avevi già scritto
            alpha0_deg = 10.0

            idx_beta1 = DOF_NAMES_ALL.index("DOF_BETA1")  # o DOF_BETA1_GEOM_ se usi quest'ultimo
            idx_beta2 = DOF_NAMES_ALL.index("DOF_BETA2")  # o DOF_BETA2_GEOM_

            beta1_deg = best_dof[idx_beta1]
            beta2_deg = best_dof[idx_beta2]

            # === CALCOLI ===
            # Deflessione = alpha_0 - alpha_ex (presumo in gradi)
            deflessione = alpha0_deg - alpha_ex_deg

            # Curvatura Camber = beta1 - beta2
            curvatura_camber = beta1_deg - beta2_deg

            # === STAMPA NEL PPTX - COLONNA SINISTRA ===
            # (baseline o task1, il codice per aggiungere è lo stesso)

            # Stampa preesistente alpha0
            p = tf_sin_best.add_paragraph()
            p.text = f"ALPHA_0 = {alpha0_deg:.3f}°"
            p.level = 0

            # Stampa Deflessione
            p_defl = tf_sin_best.add_paragraph()
            p_defl.text = f"Deflessione (ε) = {deflessione:.3f}°"
            p_defl.level = 0

            # Stampa Curvatura
            p_curv = tf_sin_best.add_paragraph()
            p_curv.text = f"Curvatura camber (θ) = {curvatura_camber:.3f}°"
            p_curv.level = 0


        except Exception as e:
            p = tf_sin_best.add_paragraph()
            p.text = f"[Info extra] Errore calcolo deflessione/curvatura: {e}"
            p.level = 0

        for i, (nome, val) in enumerate(zip(OF_NAMES, best_of)):
            p_of = tf_des_best.add_paragraph()

            # Poiché gli OF sono output, tecnicamente cambiano tutti.
            # Verifichiamo se c'è stata una variazione significativa.
            if start_of is not None:
                if abs(val - start_of[i]) > 1e-6:
                    p_of.text = f"{nome} = {val:.6f} da {start_of[i]:.6f}"
                else:
                    p_of.text = f"{nome} = {val:.6f}"
                p_of.level = 0
            else:
                p_of.text = f"{nome} = {val:.6f}"
                p_of.level = 0


def aggiungi_smith(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    if slide.shapes.title:
        slide.shapes.title.text = "Smith Diagram Action - Axial exit"
    slide.shapes.add_picture("smith_diagram_action_assiale.png", Inches(2.2), Inches(1.2),
                             height=Inches(6.2))

    # Slide 2: Smith Diagram - Action Total to Total
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    if slide.shapes.title:
        slide.shapes.title.text = "Smith Diagram Action - Total to Total"
    slide.shapes.add_picture("smith_diagram_action_total_to_total.png", Inches(2.2), Inches(1.2),
                             height=Inches(6.2))

    # Slide 3: Smith Diagram - Reaction Total to Total
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    if slide.shapes.title:
        slide.shapes.title.text = "Smith Diagram Reaction - Total to Total"
    slide.shapes.add_picture("smith_diagram_reaction_total_to_total.png", Inches(2.2), Inches(1.2),
                             height=Inches(6.2))