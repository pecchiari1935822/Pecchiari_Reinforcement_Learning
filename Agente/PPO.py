import numpy as np
import re
import pandas as pd
from pathlib import Path
import matplotlib.pyplot as plt
import matplotlib.gridspec as gridspec
from stable_baselines3 import PPO
from stable_baselines3.common.env_checker import check_env
from stable_baselines3.common.callbacks import BaseCallback, CheckpointCallback
from stable_baselines3.common.monitor import Monitor
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import time
import shutil
import glob
from Ambiente.Ambiente_claude_senza_keras import (
    SURROGATE_MODEL_PATH, SCALER_PATH,
    DOF_NAMES_ALL, OF_NAMES, ACTIVE_DOF_INDICES,
    load_surrogate, BladeOptimEnv, DOF_BOUNDS_ALL)
from Agente.Set_input_param import combinazioni_da_testare, learning_rate, n_steps, ROW_INDEX, PPO_PARAMS, TOTAL_TIMESTEPS



# ============================================================
# CALLBACK — traccia CSI e Score per episodio
# ============================================================

class BladeCallback(BaseCallback):
    """
    Traccia per ogni episodio:
      - CSI finale (ultimo step dell'episodio)
      - Score = reward cumulativa dell'episodio (come nel paper)
      - DOF e OF del miglior risultato trovato in tutto il training
    """

    def __init__(self, verbose=0):
        super().__init__(verbose)
        self.ep_rewards = []
        self.ep_lengths = []
        self.timesteps = []

        # Early stopping
        self.patience = 4000  # numero di step senza miglioramento prima di fermare
        self.steps_senza_miglioramenti = 0

        # Metriche episodio
        self.episode_csi = []
        self.episode_scores = []
        self._current_score = 0.0
        self.best_csi_ep = np.inf
        self.best_csi = np.inf
        self.best_dof = None
        self.best_of = None

        self.best_dof_ep = None
        self.episode_best_dofs = []

        # Metriche PPO per aggiornamento
        self.metrics = {
            "explained_variance": [],
            "entropy_loss": [],
            "std": [],
            "approx_kl": [],
            "clip_fraction": [],
            "value_loss": [],
            "policy_gradient_loss": [],
            "ep_rew_mean":[],
        }
        self.n_episodes = 0
        self.metrics_episodes = []
        self.metrics_timesteps = []

    # smma le ricompense ottenute in ogni step (20) per ottenere lo score totale dell'episodio
    def _on_step(self) -> bool:
        rewards = self.locals.get("rewards", [])
        infos = self.locals.get("infos", [])
        dones = self.locals.get("dones", [])

        self.steps_senza_miglioramenti += 1

        if len(self.model.ep_info_buffer) > 0:
            # ep_info_buffer contiene info sugli episodi completati

            # Estrai ep_rew_mean (media ultimi 100 episodi)
            ep_rewards = [ep_info['r'] for ep_info in self.model.ep_info_buffer]
            if len(ep_rewards) > 0:
                ep_rew_mean = np.mean(ep_rewards)

                # Salva
                self.ep_rewards.append(ep_rew_mean)
                self.timesteps.append(self.num_timesteps)

                # Debug
                if self.num_timesteps % 100 == 0:
                    print(f"Step {self.num_timesteps}: ep_rew_mean = {ep_rew_mean:.4f}")

        for reward, info, done in zip(rewards, infos, dones):
            csi_step = info.get("csi", None)
            dof_step = info.get("dof_full", None)
            of_step = info.get("of", None)
            if csi_step is not None:
                if csi_step < self.best_csi_ep:
                    self.best_csi_ep = csi_step
                    self.best_dof_ep = dof_step.copy() if dof_step is not None else None

                if csi_step < self.best_csi:
                    self.best_csi = csi_step
                    self.best_dof = dof_step.copy() if dof_step is not None else None
                    self.best_of = of_step.copy() if of_step is not None else None
                    self.steps_senza_miglioramenti = 0
                    self.logger.record("custom/best_CSI", self.best_csi)
            self._current_score += reward



            if done:
                self.episode_csi.append(self.best_csi_ep)  # Ora salvi il meglio episodio!
                self.episode_scores.append(self._current_score)
                self.episode_best_dofs.append(self.best_dof_ep)
                self._current_score = 0.0
                self.n_episodes += 1

                self.logger.record("custom/CSI", self.best_csi_ep)
                self.logger.record("custom/Score", self.episode_scores[-1])

                if self.best_csi_ep < self.best_csi:
                    self.best_csi = self.best_csi_ep
                    self.steps_senza_miglioramenti = 0
                    self.best_dof = info.get("dof_full", None)
                    self.best_of = info.get("of", None)
                    self.logger.record("custom/best_CSI", self.best_csi)

                self.best_csi_ep = np.inf  # Reset per il prossimo episodio
                self.best_dof_ep = None


        if self.steps_senza_miglioramenti >= self.patience:
                print(f"\n[Early Stopping] Nessun miglioramento del CSI per {self.patience} step. Interruzione.")
                return False  # Questo ferma il PPO!

        return True

    def _on_rollout_end(self) -> None:
        """
        Chiamato dopo ogni raccolta di n_steps.
        Legge le metriche PPO dal logger interno di SB3.
        """
        log = self.model.logger.name_to_value

        chiavi_sb3 = {
            "train/explained_variance": "explained_variance",
            "train/entropy_loss": "entropy_loss",
            "train/std": "std",
            "train/approx_kl": "approx_kl",
            "train/clip_fraction": "clip_fraction",
            "train/value_loss": "value_loss",
            "train/policy_gradient_loss": "policy_gradient_loss",

        }

        valori = {
            k_nostra: log[k_sb3]
            for k_sb3, k_nostra in chiavi_sb3.items()
            if k_sb3 in log
        }

        if len(self.model.ep_info_buffer) > 0:
            ep_rewards = [ep_info['r'] for ep_info in self.model.ep_info_buffer]
            if ep_rewards:
                valori["ep_rew_mean"] = float(np.mean(ep_rewards))

        if valori:
            for chiave, valore in valori.items():
                self.metrics[chiave].append(float(valore))
            self.metrics_timesteps.append(self.num_timesteps)
            self.metrics_episodes.append(self.n_episodes)

        return True


# ============================================================
# TRAINING
# ============================================================

def train(surrogate_path=SURROGATE_MODEL_PATH,
          scaler_path=SCALER_PATH,
          start_dof=None,
          learning_rate=None, n_steps=None, batch_size=None, ROW_INDEX=None):
    print("=" * 60)
    print("  PPO Blade Optimization — Stable Baselines3")
    print(f"  DOF attivi: {[DOF_NAMES_ALL[i] for i in ACTIVE_DOF_INDICES]}")
    print(f"  Reward: minimizza CSI (OF_CSI_OP_01)")
    print("=" * 60)

    surrogate = load_surrogate(surrogate_path, scaler_path)
    assert callable(surrogate), "load_surrogate ha restituito None"

    # check_env PRIMA di Monitor (evita bug NoneType)
    print("\n  Verifica ambiente...")
    env_raw = BladeOptimEnv(surrogate, start_dof=start_dof)
    check_env(env_raw, warn=True)
    print("  Ambiente OK.\n")

    if start_dof is None:
        import re
        active_names = [DOF_NAMES_ALL[i] for i in ACTIVE_DOF_INDICES]
        safe_names = [re.sub(r'[^0-9A-Za-z]+', '', n) for n in active_names]
        active_tag = "_".join(safe_names) if safe_names else "ALL"

        # --- percorsi dinamici per checkpoint, log e monitor ---
        checkpoint_dir = f"./ppo_blade_generale_checkpoints_{active_tag}/"
        log_dir = f"./ppo_blade_generale_logs_{active_tag}/"
        monitor_file = f"./ppo_blade_generale_monitor_{active_tag}"
    else:
        import re
        active_names = [DOF_NAMES_ALL[i] for i in ACTIVE_DOF_INDICES]
        safe_names = [re.sub(r'[^0-9A-Za-z]+', '', n) for n in active_names]
        active_tag = "_".join(safe_names) if safe_names else "ALL"

        # --- percorsi dinamici per checkpoint, log e monitor ---
        checkpoint_dir = f"./ppo_blade_start_profile_checkpoints_{active_tag}/"
        log_dir = f"./ppo_blade_start_profile_logs_{active_tag}/"
        monitor_file = f"./ppo_blade_start_profile_monitor_{active_tag}"

        # Crea le directory se non esistono
        import os
        os.makedirs(checkpoint_dir, exist_ok=True)
        os.makedirs(log_dir, exist_ok=True)


    env = Monitor(env_raw, filename="./ppo_blade_monitor")

    cb_blade = BladeCallback(verbose=0)
    cb_ckpt = CheckpointCallback(
        save_freq=10000, save_path="./ppo_blade_checkpoints/",
        name_prefix="ppo_blade", verbose=1,
    )

    model = PPO(
        policy="MlpPolicy", env=env,
        policy_kwargs=dict(net_arch=dict(pi=[128, 128], vf=[128, 128])),
        tensorboard_log=log_dir,
        seed=42,
        learning_rate=learning_rate,
        n_steps=n_steps,
        batch_size=batch_size,
        **PPO_PARAMS
    )

    print(f"  Addestramento: {TOTAL_TIMESTEPS:,} step")
    print("  (tensorboard --logdir ./ppo_blade_logs/)\n")

    # --- costruisci tag identificativo basato sui nomi DOF attivi ---
    if start_dof is None:
        active_names = [DOF_NAMES_ALL[i] for i in ACTIVE_DOF_INDICES]
        safe_names = [re.sub(r'[^0-9A-Za-z]+', '', n) for n in active_names]
        active_tag = "_".join(safe_names) if safe_names else "ALL"
        model_basename = f"ppo_blade_start_profile_{active_tag}_lr{learning_rate}_nsteps{n_steps}_riga{ROW_INDEX}"
        model_path = model_basename  # SB3 aggiunge .zip se serve

        print(f"  Addestramento: {TOTAL_TIMESTEPS:,} step")
        print(f"  Modello di output: {model_path}.zip")
        print("  (tensorboard --logdir ./ppo_blade_logs/)\n")

        start_time = time.time()

        model.learn(
            total_timesteps=TOTAL_TIMESTEPS,
            callback=[cb_blade, cb_ckpt],
            progress_bar=True,
        )

        model.save(model_path)
        print(f"\n  Modello salvato: {model_path}.zip")

    else:
        active_names = [DOF_NAMES_ALL[i] for i in ACTIVE_DOF_INDICES]
        safe_names = [re.sub(r'[^0-9A-Za-z]+', '', n) for n in active_names]
        active_tag = "_".join(safe_names) if safe_names else "ALL"
        model_basename = f"ppo_blade_start_profile_{active_tag}_lr{learning_rate}_nsteps{n_steps}_riga{ROW_INDEX}"
        model_path = model_basename  # SB3 aggiunge .zip se serve

        print(f"  Addestramento: {TOTAL_TIMESTEPS:,} step")
        print(f"  Modello di output: {model_path}.zip")
        print("  (tensorboard --logdir ./ppo_blade_logs/)\n")

        start_time = time.time()

        model.learn(
            total_timesteps=TOTAL_TIMESTEPS,
            callback=[cb_blade, cb_ckpt],
            progress_bar=True,
        )

        model.save(model_path)
        print(f"\n  Modello salvato: {model_path}.zip")

    end_time = time.time()
    training_time = end_time - start_time

    _print_results(cb_blade)
    _plot_results(cb_blade, learning_rate, n_steps, training_time=training_time)
    _plot_training_metrics_actor(cb_blade, learning_rate, n_steps)
    _plot_training_metrics_critic(cb_blade, learning_rate, n_steps)
    _plot_dof_evolution(cb_blade, learning_rate, n_steps, start_dof=start_dof)
    _plot_dof_evolution_barre(cb_blade, learning_rate, n_steps, start_dof=start_dof)

    env.close()

    return model, cb_blade.best_dof, cb_blade.best_of, cb_blade.best_csi, model_path


# ============================================================
# GRAFICI — stile paper Dussauge 2023
# ============================================================

def _plot_results(cb: BladeCallback, lr, n_step, save_path="plot_results.png", training_time=None):
    if not cb.episode_csi:
        print("  Nessun dato da plottare.")
        return

    time_str = ""
    if training_time is not None:
        mins = int(training_time // 60)
        secs = int(training_time % 60)
        time_str = f"  |  Durata: {mins}m {secs}s"

    csi_arr = np.array(cb.episode_csi)
    score_arr = np.array(cb.episode_scores)
    n_ep = len(csi_arr)
    ep_axis = np.arange(n_ep)
    W = min(40, max(n_ep // 5, 2))

    def moving_avg(arr, w):
        if len(arr) < w:
            return ep_axis, arr
        ma = np.convolve(arr, np.ones(w) / w, mode='valid')
        return np.arange(w - 1, len(arr)), ma

    fig = plt.figure(figsize=(16, 10))
    fig.suptitle(
        f"PPO Blade Opt. — DOF attivi: {[DOF_NAMES_ALL[i] for i in ACTIVE_DOF_INDICES]}\n"
        f"Total steps={TOTAL_TIMESTEPS:,}  "
        f"n_steps={n_step}" 
        f" Learning_Rate={lr}"
        f"{time_str}",
        fontsize=16
    )
    gs = gridspec.GridSpec(2, 3, figure=fig, hspace=0.4, wspace=0.35)

    # ── 1: CSI per episodio ──
    ax1 = fig.add_subplot(gs[0, 0])
    ax1.scatter(ep_axis, csi_arr, s=8, alpha=0.4, color='orange', label='CSI')
    x_ma, ma = moving_avg(csi_arr, W)
    ax1.plot(x_ma, ma, color='steelblue', lw=1.5, label=f'<CSI> ({W}ep)')
    ax1.plot(ep_axis, np.minimum.accumulate(csi_arr),
             color='green', lw=2, label='Min CSI')
    ax1.axhline(cb.best_csi, color='green', ls='--', lw=1)
    ax1.set_xlabel('Episodio', fontsize=12)
    ax1.set_ylabel('CSI', fontsize=12)
    ax1.set_title('CSI per episodio', fontsize=12)
    ax1.legend(fontsize=9, loc='upper right')
    ax1.grid(alpha=0.3)
    ax1.tick_params(labelsize=11)

    # ── 2: Score per episodio ──
    ax2 = fig.add_subplot(gs[0, 1])
    ax2.scatter(ep_axis, score_arr, s=8, alpha=0.4, color='orange', label='Score')
    x_ma2, ma2 = moving_avg(score_arr, W)
    ax2.plot(x_ma2, ma2, color='steelblue', lw=1.5, label=f'<Score> ({W}ep)')
    ax2.plot(ep_axis, np.maximum.accumulate(score_arr),
             color='green', lw=2, label='Max Score')
    ax2.set_xlabel('Episodio', fontsize=12)
    ax2.set_ylabel('Score (reward cumulativa)', fontsize=12)
    ax2.set_title('Score per episodio', fontsize=12)
    ax2.legend(fontsize=9)
    ax2.grid(alpha=0.3)
    ax2.tick_params(labelsize=11)

    # ── 3: CSI minimo progressivo ──
    ax3 = fig.add_subplot(gs[0, 2])
    ax3.plot(ep_axis, np.minimum.accumulate(csi_arr), color='green', lw=2)
    ax3.annotate(f'Min CSI = {cb.best_csi:.5f}',
                 xy=(0.05, 0.05), xycoords='axes fraction', fontsize=10,
                 color='green',
                 bbox=dict(boxstyle='round,pad=0.3', fc='white', ec='green', alpha=0.8))
    ax3.set_xlabel('Episodio', fontsize=12)
    ax3.set_ylabel('Miglior CSI trovato', fontsize=12)
    ax3.set_title('CSI minimo progressivo', fontsize=12)
    ax3.grid(alpha=0.3)
    ax3.tick_params(labelsize=11)

    plt.tight_layout(rect=[0, 0, 1, 0.93])
    plt.savefig(save_path, bbox_inches='tight', dpi=150)
    plt.close()

    return save_path


def _plot_dof_evolution(cb: BladeCallback, lr, n_step, start_dof = None, save_path="plot_dof_evolution.png"):
    """
    Crea un grafico per ogni DOF mostrando l'andamento del miglior valore trovato
    in ogni episodio per tutta la durata dell'addestramento, colorando i punti in
    base al numero dell'episodio.
    """
    if not cb.episode_best_dofs:
        print("  Nessun dato dei DOF da plottare.")
        return

    # Trasforma la lista di array in una matrice (n_episodi, 7_dof)
    dof_data = np.array(cb.episode_best_dofs)
    n_episodes = dof_data.shape[0]
    ep_axis = np.arange(n_episodes)

    best_ep_idx = int(np.argmin(cb.episode_csi))

    np.random.seed(42)  # Fissa il seed così i colori non cambiano a ogni run
    colori_per_episodio = np.random.uniform(0.1,0.75,size=(n_episodes, 3))

    # Creiamo una griglia 4x2 standard (senza sharex)
    fig, axes = plt.subplots(4, 2, figsize=(15, 12))
    fig.suptitle(
        f"Evoluzione dei DOF (Miglior profilo per Episodio)\n"
        f"Total steps={TOTAL_TIMESTEPS:,}  n_steps={n_step}  Learning_Rate={lr}",
        fontsize=16
    )
    axes_flat = axes.flatten()

    for i in range(7):
        ax = axes_flat[i]
        y_vals = dof_data[:, i]

        if i in ACTIVE_DOF_INDICES:
            title_suffix = " (Attivo)"
            ma_color = "black"

            # Sfumatura di colore: c=ep_axis mappa il colore sul numero dell'episodio
            # cmap='viridis' va da viola scuro (inizio), a verde (centro), a giallo (fine)
            ax.scatter(ep_axis, y_vals, c=colori_per_episodio, s=15, alpha=0.8)

            # Elemento "fantasma" per mostrare l'etichetta nella legenda
            # (altrimenti scatter con c= array non crea un'etichetta semplice)
            ax.plot([], [], 'o', color='mediumseagreen', markersize=5, label="Miglior DOF/Ep")
        else:
            title_suffix = " (Fisso)"
            ma_color = "black"
            ax.scatter(ep_axis, y_vals, color='gray', s=15, alpha=0.5, label="Fisso")

        # Aggiungi una media mobile per vedere meglio il trend
        '''if n_episodes > 10:
            w = max(5, n_episodes // 20)
            ma = np.convolve(y_vals, np.ones(w) / w, mode='valid')
            ax.plot(np.arange(w - 1, n_episodes), ma, color=ma_color, lw=2, label=f"Media mobile ({w} ep)")'''

        if start_dof is not None:
            start_val = start_dof[i]
            # Mettiamo il cerchio all'episodio 0 (inizio training)
            ax.plot(0, start_val, marker='o', color='cyan', markeredgecolor='black',
                    markersize=10, linestyle='None', zorder=5, label="Partenza")

        if cb.best_dof is not None:
            best_val = cb.best_dof[i]
            # Usiamo zorder=5 per assicurarci che la X venga disegnata SOPRA le barre e le linee
            ax.plot(best_ep_idx, best_val, marker='X', color='red', markeredgecolor='black',
                    markersize=12, linestyle='None', zorder=5, label="Miglior Assoluto")

        # Disegna i limiti fisici (bounds)
        dof_min, dof_max = DOF_BOUNDS_ALL[i]
        ax.axhline(dof_min, color='red', ls='--', lw=1, alpha=0.5)
        ax.axhline(dof_max, color='red', ls='--', lw=1, alpha=0.5)

        ax.set_title(DOF_NAMES_ALL[i] + title_suffix, fontsize=12)
        ax.set_xlabel("Episodio", fontsize=10)
        ax.set_ylabel("Valore", fontsize=10)
        ax.grid(alpha=0.3)

        if i in ACTIVE_DOF_INDICES:
            ax.legend(fontsize=8, loc='best')

    # Nascondi l'ottavo grafico vuoto
    axes_flat[7].set_visible(False)

    plt.tight_layout(rect=[0, 0, 1, 0.95])
    plt.savefig(save_path, bbox_inches='tight', dpi=150)
    plt.close()

    return save_path


def _plot_dof_evolution_barre(cb: BladeCallback, lr, n_step,start_dof = None, save_path="plot_dof_evolution_barre.png"):
    """
    Crea un grafico per ogni DOF mostrando l'andamento del miglior valore trovato
    in ogni episodio. Usa barre verticali colorate casualmente per ogni episodio.
    """
    if not cb.episode_best_dofs:
        print("  Nessun dato dei DOF da plottare.")
        return

    # Trasforma la lista di array in una matrice (n_episodi, 7_dof)
    dof_data = np.array(cb.episode_best_dofs)
    n_episodes = dof_data.shape[0]
    ep_axis = np.arange(n_episodes)

    best_ep_idx = int(np.argmin(cb.episode_csi))

    # Generiamo un colore RGB casuale per OGNI episodio
    np.random.seed(42)
    colori_per_episodio = np.random.uniform(0.1,0.75,size=(n_episodes, 3))

    # Creiamo una griglia 4x2 standard
    fig, axes = plt.subplots(4, 2, figsize=(15, 12))
    fig.suptitle(
        f"Evoluzione dei DOF (Miglior profilo per Episodio)\n"
        f"Total steps={TOTAL_TIMESTEPS:,}  n_steps={n_step}  Learning_Rate={lr}",
        fontsize=16
    )
    axes_flat = axes.flatten()

    for i in range(7):
        ax = axes_flat[i]
        y_vals = dof_data[:, i]

        if i in ACTIVE_DOF_INDICES:
            title_suffix = " (Attivo)"
            ma_color = "black"

            # Disegna le BARRE colorate
            ax.bar(ep_axis, y_vals, color="dodgerblue", width=1.0, alpha=0.9)

            # Elemento fantasma per la legenda (barra vuota)
            ax.bar([-10], [0], color='gray', label="Miglior DOF/Ep")
        else:
            title_suffix = " (Fisso)"
            ma_color = "black"
            ax.bar(ep_axis, y_vals, color='gray', width=1.0, alpha=0.5, label="Fisso")

        # Aggiungi una media mobile per vedere meglio il trend
        '''if n_episodes > 10:
            w = max(5, n_episodes // 20)
            ma = np.convolve(y_vals, np.ones(w) / w, mode='valid')
            ax.plot(np.arange(w - 1, n_episodes), ma, color=ma_color, lw=2, label=f"Media mobile ({w} ep)")'''

        if start_dof is not None:
            start_val = start_dof[i]
            # Mettiamo il cerchio all'episodio 0 (inizio training)
            ax.plot(0, start_val, marker='o', color='cyan', markeredgecolor='black',
                    markersize=10, linestyle='None', zorder=5, label="Partenza")

        if cb.best_dof is not None:
            best_val = cb.best_dof[i]
            # Usiamo zorder=5 per assicurarci che la X venga disegnata SOPRA le barre e le linee
            ax.plot(best_ep_idx, best_val, marker='X', color='red', markeredgecolor='black',
                    markersize=12, linestyle='None', zorder=5, label="Miglior Assoluto")

        # Disegna i limiti fisici (bounds)
        dof_min, dof_max = DOF_BOUNDS_ALL[i]
        ax.axhline(dof_min, color='red', ls='--', lw=1, alpha=0.5)
        ax.axhline(dof_max, color='red', ls='--', lw=1, alpha=0.5)

        # --- IMPORTANTE: FORZA LO ZOOM DELL'ASSE Y ---
        # Evita che le barre partano forzatamente da 0 rovinando la scala
        padding = (dof_max - dof_min) * 0.1
        if padding == 0: padding = 0.1
        ax.set_ylim(dof_min - padding, dof_max + padding)

        ax.set_title(DOF_NAMES_ALL[i] + title_suffix, fontsize=12)
        ax.set_xlabel("Episodio", fontsize=10)
        ax.set_ylabel("Valore", fontsize=10)
        ax.grid(alpha=0.3)
        ax.set_xlim(left=-1, right=n_episodes)

        if i in ACTIVE_DOF_INDICES:
            ax.legend(fontsize=8, loc='best')

    # Nascondi l'ottavo grafico vuoto
    axes_flat[7].set_visible(False)

    plt.tight_layout(rect=[0, 0, 1, 0.95])
    plt.savefig(save_path, bbox_inches='tight', dpi=150)
    plt.close()

    return save_path


def _plot_training_metrics_actor(cb: BladeCallback, lr, n_step, save_path="plot_metrics_actor.png"):

    if not cb.metrics_episodes:
        print("  Nessuna metrica PPO disponibile (training troppo corto).")
        return

    ts = np.array(cb.metrics_episodes)

    # Configurazione: (chiave, titolo, colore, linea riferimento, label ref)
    config = [("ep_rew_mean", "Episode Reward Mean", "navy", None, None),
            ("policy_gradient_loss", "Policy Gradient Loss (L_CLIP Actor)", "darkslategray", 0.0, "zero"),
            ("entropy_loss", "Entropy H[π] (esplorazione Actor)", "darkorange", None, None),
            ("std", "Std Policy (deviazione standard azioni Actor)", "mediumpurple", None, None),
            ("approx_kl", "Approx KL Divergence (cambio policy per update Actor)", "crimson", 0.02, "soglia 0.02"),
            ("clip_fraction", "Clip Fraction ( % azioni clippate)", "teal", 0.1, "soglia 0.1")
    ]

    fig, axes = plt.subplots(3, 3, figsize=(16, 11))
    fig.suptitle(
        f"Metriche interne Actor — DOF: {[DOF_NAMES_ALL[i] for i in ACTIVE_DOF_INDICES]}\n"
        f"Total steps={TOTAL_TIMESTEPS:,}  n_steps={n_step}  Learning_Rate={lr} ",
        fontsize=16
    )
    axes_flat = axes.flatten()

    for idx, (chiave, titolo, colore, ref_val, ref_label) in enumerate(config):
        ax = axes_flat[idx]
        valori = cb.metrics.get(chiave, [])

        if not valori:
            ax.text(0.5, 0.5, "Nessun dato", ha="center", va="center",
                    transform=ax.transAxes, color="gray")
            ax.set_title(titolo, fontsize=12)
            continue

        if chiave == "entropy_loss":
            vals = -np.array(valori)
        else:
            vals = np.array(valori)
        ax.plot(ts[:len(vals)], vals, color=colore, lw=1.5, alpha=0.85)

        # Media mobile per leggibilità
        if len(vals) >= 5:
            w = max(3, len(vals) // 10)
            ma = np.convolve(vals, np.ones(w) / w, mode="valid")
            ax.plot(ts[w - 1:len(vals)], ma,
                    color=colore, lw=2.5, alpha=0.5,
                    linestyle="--", label=f"media {w} update")

        # Linea di riferimento (se presente)
        if ref_val is not None:
            ax.axhline(ref_val, color="green", ls=":", lw=1.2,
                       alpha=0.7, label=ref_label)

        # Valore finale annotato
        ax.annotate(
            f"finale: {vals[-1]:.4f}",
            xy=(ts[len(vals) - 1], vals[-1]),
            xytext=(-60, 8), textcoords="offset points",
            fontsize=10, color=colore,
            arrowprops=dict(arrowstyle="->", color=colore, lw=0.8)
        )

        ax.set_xlabel("Episodio", fontsize=10)
        ax.set_title(titolo, fontsize=11, fontweight="bold")
        ax.legend(fontsize=10)
        ax.grid(alpha=0.25)
        ax.tick_params(labelsize=11)

    # Disabilita gli ultimi 2 assi (griglia 3x3, metriche sono 7)
    for idx in range(len(config), len(axes_flat)):
        axes_flat[idx].set_visible(False)

    plt.tight_layout(rect=[0, 0, 1, 0.95])
    plt.savefig(save_path, bbox_inches='tight', dpi=150)
    plt.close()

    return save_path



def _plot_training_metrics_critic(cb: BladeCallback,lr , n_step, save_path="plot_metrics_critic.png"):
    if not cb.metrics_episodes:
        print("  Nessuna metrica PPO disponibile (training troppo corto).")
        return

    ts = np.array(cb.metrics_episodes)

    # Configurazione: (chiave, titolo, colore, linea riferimento, label ref)
    config = [("explained_variance", "Explained Variance (Critic)", "steelblue", 1.0, "ottimo=1.0"),
            ("value_loss", "Value Loss (errore Critic sui Returns)", "sienna", None, None),
              ]

    fig, axes = plt.subplots(1, 2, figsize=(16, 6))
    fig.suptitle(
        f"Metriche interne Actor — DOF: {[DOF_NAMES_ALL[i] for i in ACTIVE_DOF_INDICES]}\n"
        f"Total steps={TOTAL_TIMESTEPS:,}  n_steps={n_step}  Learning_Rate={lr} ",
        fontsize=16
    )
    axes_flat = axes.flatten()

    for idx, (chiave, titolo, colore, ref_val, ref_label) in enumerate(config):
        ax = axes_flat[idx]
        valori = cb.metrics.get(chiave, [])

        if not valori:
            ax.text(0.5, 0.5, "Nessun dato", ha="center", va="center",
                    transform=ax.transAxes, color="gray")
            ax.set_title(titolo, fontsize=12)
            continue

        if chiave == "entropy_loss":
            vals = -np.array(valori)
        else:
            vals = np.array(valori)
        ax.plot(ts[:len(vals)], vals, color=colore, lw=1.5, alpha=0.85)

        # Media mobile per leggibilità
        if len(vals) >= 5:
            w = max(3, len(vals) // 10)
            ma = np.convolve(vals, np.ones(w) / w, mode="valid")
            ax.plot(ts[w - 1:len(vals)], ma,
                    color=colore, lw=2.5, alpha=0.5,
                    linestyle="--", label=f"media {w} update")

        # Linea di riferimento (se presente)
        if ref_val is not None:
            ax.axhline(ref_val, color="green", ls=":", lw=1.2,
                       alpha=0.7, label=ref_label)

        # Valore finale annotato
        ax.annotate(
            f"finale: {vals[-1]:.4f}",
            xy=(ts[len(vals) - 1], vals[-1]),
            xytext=(-60, 8), textcoords="offset points",
            fontsize=10, color=colore,
            arrowprops=dict(arrowstyle="->", color=colore, lw=0.8)
        )

        ax.set_xlabel("Episodio", fontsize=10)
        ax.set_title(titolo, fontsize=11, fontweight="bold")
        ax.legend(fontsize=10)
        ax.grid(alpha=0.25)
        ax.tick_params(labelsize=11)

    # Disabilita gli ultimi 2 assi (griglia 3x3, metriche sono 7)
    for idx in range(len(config), len(axes_flat)):
        axes_flat[idx].set_visible(False)

    plt.tight_layout(rect=[0, 0, 1, 0.95])
    plt.savefig(save_path, bbox_inches='tight', dpi=150)
    plt.close()

    return save_path

# ============================================================
# UTILITY
# ============================================================

def _print_results(cb: BladeCallback):
    print(f"\n  Miglior CSI durante training: {cb.best_csi:.6f}")
    if cb.best_dof is not None:
        print("\n  DOF ottimali che mi danno il miglior CSI (* = modificato):")
        for i, (name, val) in enumerate(zip(DOF_NAMES_ALL, cb.best_dof)):
            m = " *" if i in ACTIVE_DOF_INDICES else "  "
            print(f"    {name:<28}{m}: {val:.6f}")
        print("\n  OF corrispondenti ai DOF ottimali:")
        for name, val in zip(OF_NAMES, cb.best_of):
            print(f"    {name:<30}: {val:.6f}")

    print("\n  Questo è il profilo che ha il miglior CSI Durante il training")


# ============================================================
# Presentazione PowerPoint — aggiunge slide per ogni iterazione di training con parametri e grafici
# ============================================================

def aggiungi_slide_iterazione(prs, parametri, img_paths, row_idx, lr, best_dof, best_of,start_dof, start_of):
    """
    Aggiunge 4 slide alla presentazione: 1 di testo e 3 di immagini.
    """
    slide_layout_testo = prs.slide_layouts[3]
    slide_testo = prs.slides.add_slide(slide_layout_testo)

    if slide_testo.shapes.title:
        title_shape = slide_testo.shapes.title
        title_shape.text = f"Punto di partenza simulazione con LR = {lr}"


    # 1. Slide di testo con i parametri

    if len(slide_testo.placeholders) > 1:
        corpo_sinistro = slide_testo.placeholders[1]
        tf_sinistra = corpo_sinistro.text_frame
        tf_sinistra.clear()  # Pulisce il testo di default
        for key, value in parametri.items():
            if isinstance(value, list):
                # 1. Aggiungi la voce principale (es. "DOF attivi modificati:")
                p_main = tf_sinistra.add_paragraph()
                p_main.text = f"{key}:"
                p_main.level = 0

                # 2. Aggiungi i vari elementi della lista indentati (level = 1)
                for item in value:
                    p_sub = tf_sinistra.add_paragraph()
                    p_sub.text = str(item)
                    p_sub.level = 1  # QUESTO CREA L'INDENTAZIONE!

            # Se è un valore normale (testo o numero), lo mette come punto principale
            else:
                p = tf_sinistra.add_paragraph()
                p.level = 0

                if key.strip() == "":
                    p.text = ""  # riga vuota
                else:
                    p.text = f"{key} = {value}"



        # COLONNA DESTRA (Limiti dei DOF)
        # =======================================================
        # Posizioniamo una nuova casella di testo nella metà destra (da 7 pollici in poi)
        if len(slide_testo.placeholders) > 2:
            corpo_destro = slide_testo.placeholders[2]
            tf_destra = corpo_destro.text_frame
            tf_destra.clear()  # Cancella la scritta "Fare clic per inserire testo"

            p_title_destra = tf_destra.add_paragraph()
            p_title_destra.text = "Limiti operativi (DOF Bounds):"
            p_title_destra.level = 0

            # Aggiungiamo l'elenco dei limiti
            for nome_dof, bounds in zip(DOF_NAMES_ALL, DOF_BOUNDS_ALL):
                p_bound = tf_destra.add_paragraph()
                p_bound.text = f"{nome_dof} = [{bounds[0]}, {bounds[1]}]"


    # 2. Slide per ogni immagine (3 immagini)
    try:
        slide_layout_img = prs.slide_layouts[4]  # Prova con Solo Titolo
    except:
        slide_layout_img = prs.slide_layouts[1]  # Fallback

    titoli_immagini = ["Risultati","Andamento Miglior DOF per episodio", "Andamento Miglior DOF per episodio", "Metriche Attore", "Metriche Critico"]  # Layout Solo Titolo o Vuota
    for idx, img_path in enumerate(img_paths):
        if not os.path.exists(img_path):
            continue

        slide_img = prs.slides.add_slide(slide_layout_img)

        if slide_img.shapes.title:
            slide_img.shapes.title.text = f"{titoli_immagini[idx]}"

        # Inserisci immagine centrata (adattata per 16:9)
        # Una slide 16:9 tipica è larga 13.33 pollici e alta 7.5 pollici
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(11.33)  # Lascia 1 pollice di margine per lato

        slide_img.shapes.add_picture(str(img_path), left, top, width=width)

    slide_layout_best = prs.slide_layouts[5]
    slide_best = prs.slides.add_slide(slide_layout_best)

    if slide_best.shapes.title:
        title_shape = slide_best.shapes.title
        title_shape.text = f"Miglior Profilo Ottimizzato"


    # COLONNA SINISTRA: I 7 DOF ottimali
    if len(slide_best.placeholders) > 1:
        corpo_sin_best = slide_best.placeholders[1]
        tf_sin_best = corpo_sin_best.text_frame
        tf_sin_best.clear()

        p_title_dof = tf_sin_best.add_paragraph()
        p_title_dof.text = "DOF Ottimali:"
        p_title_dof.level = 0

        for i, (nome, val) in enumerate(zip(DOF_NAMES_ALL, best_dof)):
            p_dof = tf_sin_best.add_paragraph()

            # Se il DOF era attivo, mostra anche il valore di partenza
            if i in ACTIVE_DOF_INDICES:
                p_dof.text = f"{nome} (*) = {val:.6f} da {start_dof[i]:.6f}"
            else:
                p_dof.text = f"{nome} = {val:.6f}"
            p_dof.level = 0

    # COLONNA DESTRA: I 15 OF corrispondenti
    if len(slide_best.placeholders) > 2:
        corpo_des_best = slide_best.placeholders[2]
        tf_des_best = corpo_des_best.text_frame
        tf_des_best.clear()

        p_title_of = tf_des_best.paragraphs[0]
        p_title_of.text = "OF Risultanti:"
        p_title_of.level = 0

        for i, (nome, val) in enumerate(zip(OF_NAMES, best_of)):
            p_of = tf_des_best.add_paragraph()

            # Poiché gli OF sono output, tecnicamente cambiano tutti.
            # Verifichiamo se c'è stata una variazione significativa.
            if abs(val - start_of[i]) > 1e-6:
                p_of.text = f"{nome} = {val:.6f} da {start_of[i]:.6f}"
            else:
                p_of.text = f"{nome} = {val:.6f}"
            p_of.level = 0


def pulisci_file_temporanei():
    """Elimina log, immagini temporanee e checkpoint, mantenendo solo Modelli e PPTX."""
    print("\n" + "=" * 60)
    print("  PULIZIA FILE TEMPORANEI")
    print("=" * 60)

    # 1. Elimina le immagini dei grafici
    immagini = ["plot_results.png", "plot_metrics_actor.png", "plot_metrics_critic.png","plot_dof_evolution_barre.png", "plot_dof_evolution.png"]
    for img in immagini:
        if os.path.exists(img):
            os.remove(img)
            print(f"  [X] Eliminato: {img}")

    # 2. Elimina i file CSV del Monitor di Gym
    for monitor_file in glob.glob("*monitor*.csv"):
        os.remove(monitor_file)
        print(f"  [X] Eliminato: {monitor_file}")

    # 3. Elimina le cartelle di Checkpoint e Log di TensorBoard
    # Trova tutte le cartelle che corrispondono a questi pattern
    cartelle_da_eliminare = [
        "ppo_blade_checkpoints",  # Quella hardcoded di base
    ]
    cartelle_da_eliminare.extend(glob.glob("ppo_blade_*_checkpoints_*"))
    cartelle_da_eliminare.extend(glob.glob("ppo_blade_*_logs_*"))

    for cartella in cartelle_da_eliminare:
        if os.path.exists(cartella) and os.path.isdir(cartella):
            shutil.rmtree(cartella)  # rmtree cancella la cartella e tutto il suo contenuto
            print(f"  [X] Eliminata cartella: {cartella}")

    print("=" * 60)
    print("  Pulizia completata! Conservati solo i file .zip finali e la presentazione.")
    print("=" * 60)




